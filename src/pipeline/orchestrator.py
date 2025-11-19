"""Pipeline orchestrator for end-to-end CV automation."""
from typing import Dict, Any, List, Optional
from pathlib import Path
import json
from datetime import datetime
from pptx.util import Pt

from ..agents.agent1_input_validator import Agent1_InputValidator
from ..agents.agent2_template_loader import Agent2_TemplateLoader
from ..agents.agent3_semantic_extractor import Agent3_SemanticExtractor
from ..agents.agent4_mapping_engine import Agent4_MappingEngine
from ..agents.agent5_placeholder_analyzer import Agent5_PlaceholderAnalyzer
from ..agents.agent6_content_transformer import Agent6_ContentTransformer
from ..core.exceptions import CVAutomationError
from ..utils.logger import logger
from .traceability import TraceabilityManager


class CVAutomationPipeline:
    """
    End-to-end pipeline orchestrator.
    Implements REQ-PIPE-01: Sequential deterministic execution
    """

    def __init__(self):
        self.agent1 = Agent1_InputValidator()
        self.agent2 = Agent2_TemplateLoader()
        self.agent3 = Agent3_SemanticExtractor()
        self.agent4 = Agent4_MappingEngine()
        self.agent5 = Agent5_PlaceholderAnalyzer()
        self.agent6 = Agent6_ContentTransformer()

        self.traceability = TraceabilityManager()
        self.pipeline_log = []

    def execute(
        self,
        cv_file_path: str,
        template_path: str,
        output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Execute complete pipeline.

        Args:
            cv_file_path: Path to CV file
            template_path: Path to PowerPoint template
            output_path: Optional output path for generated presentation

        Returns:
            Complete pipeline result with traceability
        """
        pipeline_id = f"pipeline_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        logger.info(f"Starting pipeline execution: {pipeline_id}")

        try:
            # Stage 1: Input Validation
            logger.info("Stage 1: Validating input CV file")
            agent1_result = self.agent1.execute({"file_path": cv_file_path})

            if agent1_result.get("status") != "success":
                raise CVAutomationError(f"Agent 1 failed: {agent1_result.get('error')}")

            self.traceability.add_stage("agent1_input_validation", agent1_result)

            # Stage 2: Template Loading
            logger.info("Stage 2: Loading PowerPoint template")
            agent2_result = self.agent2.execute({"template_path": template_path})

            if agent2_result.get("status") != "success":
                raise CVAutomationError(f"Agent 2 failed: {agent2_result.get('error')}")

            self.traceability.add_stage("agent2_template_loading", agent2_result)

            # Stage 3: Semantic CV Extraction
            logger.info("Stage 3: Extracting CV data with semantic enrichment")
            agent3_result = self.agent3.execute({
                "text": agent1_result["result"]["text"],
                "file_path": cv_file_path
            })

            if agent3_result.get("status") != "success":
                raise CVAutomationError(f"Agent 3 failed: {agent3_result.get('error')}")

            self.traceability.add_stage("agent3_semantic_extraction", agent3_result)

            # Stage 4: CV-to-Template Mapping
            logger.info("Stage 4: Generating intelligent mappings")
            agent4_result = self.agent4.execute({
                "cv_data": agent3_result["result"]["cv_data_dict"],
                "placeholders": agent2_result["result"]["placeholders"]
            })

            if agent4_result.get("status") != "success":
                raise CVAutomationError(f"Agent 4 failed: {agent4_result.get('error')}")

            self.traceability.add_stage("agent4_mapping", agent4_result)

            # Stage 5: Placeholder Analysis & Rule Generation
            logger.info("Stage 5: Analyzing placeholders and generating rules")
            agent5_result = self.agent5.execute({
                "mappings": agent4_result["result"]["mappings"],
                "placeholders": agent2_result["result"]["placeholders"],
                "cv_data": agent3_result["result"]["cv_data_dict"]
            })

            if agent5_result.get("status") != "success":
                raise CVAutomationError(f"Agent 5 failed: {agent5_result.get('error')}")

            self.traceability.add_stage("agent5_placeholder_analysis", agent5_result)

            # Stage 6: Content Transformation
            logger.info("Stage 6: Transforming content with rules")

            # Prepare content mappings
            content_mappings = []
            for mapping in agent4_result["result"]["mappings"]:
                placeholder_id = mapping["placeholder_id"]
                rules = agent5_result["result"]["transformation_rules"].get(
                    placeholder_id, {}
                )

                content_mappings.append({
                    "placeholder_id": placeholder_id,
                    "source_field": mapping["source_field"],
                    "rules": rules
                })

            agent6_result = self.agent6.execute({
                "content_mappings": content_mappings,
                "cv_data": agent3_result["result"]["cv_data_dict"]
            })

            if agent6_result.get("status") != "success":
                raise CVAutomationError(f"Agent 6 failed: {agent6_result.get('error')}")

            self.traceability.add_stage("agent6_content_transformation", agent6_result)

            # Stage 7: PowerPoint Generation
            logger.info("Stage 7: Generating PowerPoint presentation")
            presentation_result = self._generate_presentation(
                agent2_result["result"]["presentation"],
                agent2_result["result"]["placeholders"],
                agent6_result["result"]["transformed_content"],
                output_path or f"outputs/generated_{pipeline_id}.pptx"
            )

            self.traceability.add_stage("presentation_generation", {
                "status": "success",
                "result": presentation_result
            })

            # Generate traceability matrix
            logger.info("Generating complete traceability matrix")
            traceability_matrix = self.traceability.generate_matrix()

            # Export traceability
            traceability_path = f"outputs/traceability/traceability_{pipeline_id}.json"
            self.traceability.export_json(traceability_path)

            traceability_excel = f"outputs/traceability/traceability_{pipeline_id}.xlsx"
            self.traceability.export_excel(traceability_excel)

            logger.info(f"Pipeline completed successfully: {pipeline_id}")

            return {
                "status": "success",
                "pipeline_id": pipeline_id,
                "output_file": presentation_result["output_path"],
                "traceability_json": traceability_path,
                "traceability_excel": traceability_excel,
                "stages": self.traceability.stages,
                "summary": {
                    "cv_file": cv_file_path,
                    "template_file": template_path,
                    "total_placeholders": len(agent2_result["result"]["placeholders"]),
                    "mapped_placeholders": len(agent4_result["result"]["mappings"]),
                    "transformations": len(agent6_result["result"]["transformed_content"]),
                    "full_compliance": agent6_result["result"]["transformation_summary"]["full_compliance"],
                    "partial_compliance": agent6_result["result"]["transformation_summary"]["partial_compliance"]
                }
            }

        except Exception as e:
            logger.error(f"Pipeline failed: {str(e)}")

            # Log failure
            self.traceability.add_stage("pipeline_failure", {
                "status": "failed",
                "error": str(e),
                "error_type": type(e).__name__
            })

            return {
                "status": "failed",
                "pipeline_id": pipeline_id,
                "error": str(e),
                "error_type": type(e).__name__,
                "stages": self.traceability.stages
            }

    def _generate_presentation(
        self,
        presentation,
        placeholders: List[Dict[str, Any]],
        transformed_content: List[Dict[str, Any]],
        output_path: str
    ) -> Dict[str, Any]:
        """Generate PowerPoint with transformed content."""

        # Create mapping of placeholder_id to transformed text
        content_map = {
            tc["placeholder_id"]: tc["transformed_text"]
            for tc in transformed_content
        }

        # Apply content to placeholders
        applied_count = 0

        for placeholder in placeholders:
            placeholder_id = placeholder["placeholder_id"]

            if placeholder_id not in content_map:
                continue

            # Find the slide and shape
            slide_idx = placeholder["slide_index"]
            shape_idx = placeholder["shape_index"]

            try:
                slide = presentation.slides[slide_idx]
                shape = slide.shapes[shape_idx]

                if shape.has_text_frame:
                    # Clear existing text
                    shape.text_frame.clear()

                    # Add new content
                    transformed_text = content_map[placeholder_id]

                    # Handle markdown formatting
                    self._apply_text_with_formatting(
                        shape.text_frame,
                        transformed_text
                    )

                    applied_count += 1

            except Exception as e:
                logger.warning(f"Failed to apply content to {placeholder_id}: {str(e)}")

        # Save presentation
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)

        logger.info(f"Saved presentation to {output_path}")

        return {
            "output_path": output_path,
            "placeholders_filled": applied_count,
            "total_placeholders": len(placeholders)
        }

    def _apply_text_with_formatting(self, text_frame, formatted_text: str):
        """Apply text with markdown formatting to text frame."""

        # Split by lines for bullet points
        lines = formatted_text.split('\n')

        for line_idx, line in enumerate(lines):
            if not line.strip():
                continue

            # Add paragraph
            if line_idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            # Handle bullet points
            if line.strip().startswith('•'):
                p.level = 0
                line = line.strip()[1:].strip()  # Remove bullet

            # Parse and apply formatting
            self._parse_and_apply_markdown(p, line)

    def _parse_and_apply_markdown(self, paragraph, text: str):
        """Parse markdown and apply formatting to paragraph."""
        # Simple markdown parser for bold and italic
        parts = []
        current = ""
        i = 0

        while i < len(text):
            if i < len(text) - 1 and text[i:i+2] == '**':
                # Bold
                if current:
                    parts.append(("normal", current))
                    current = ""

                # Find closing **
                end = text.find('**', i + 2)
                if end != -1:
                    parts.append(("bold", text[i+2:end]))
                    i = end + 2
                    continue

            elif text[i] == '*':
                # Italic
                if current:
                    parts.append(("normal", current))
                    current = ""

                # Find closing *
                end = text.find('*', i + 1)
                if end != -1:
                    parts.append(("italic", text[i+1:end]))
                    i = end + 1
                    continue

            current += text[i]
            i += 1

        if current:
            parts.append(("normal", current))

        # Apply parts to paragraph
        for style, content in parts:
            run = paragraph.add_run()
            run.text = content

            if style == "bold":
                run.font.bold = True
            elif style == "italic":
                run.font.italic = True
