"""Agent 2: PowerPoint template loading and analysis."""
from typing import Dict, Any, List
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base_agent import BaseAgent
from ..core.exceptions import TemplateError, FileIOError


class Agent2_TemplateLoader(BaseAgent):
    """
    Loads PowerPoint templates and extracts placeholder information.
    Implements REQ-A2-01: Template structure analysis
    """

    def __init__(self):
        super().__init__("agent2_template_loader", temperature=0.0)

    def process(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Load and analyze PowerPoint template.

        Args:
            input_data: {
                "template_path": str
            }

        Returns:
            {
                "template_path": str,
                "presentation": Presentation object,
                "slides": List[Dict],
                "placeholders": List[Dict],
                "metadata": Dict
            }

        Raises:
            TemplateError: If template loading or analysis fails
        """
        template_path = input_data.get("template_path")
        if not template_path:
            raise TemplateError("template_path is required")

        template_path = Path(template_path)
        if not template_path.exists():
            raise FileIOError(f"Template not found: {template_path}")

        if template_path.suffix.lower() != '.pptx':
            raise TemplateError(
                f"Invalid template format. Expected .pptx, got {template_path.suffix}"
            )

        # Load presentation
        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise TemplateError(f"Failed to load template: {str(e)}")

        # Analyze structure
        slides_data = self._analyze_slides(prs)
        placeholders_data = self._extract_placeholders(prs)

        self.logger.info(
            f"Loaded template: {len(prs.slides)} slides, {len(placeholders_data)} placeholders"
        )

        return {
            "template_path": str(template_path),
            "presentation": prs,  # Store for later use
            "slides": slides_data,
            "placeholders": placeholders_data,
            "metadata": {
                "slide_count": len(prs.slides),
                "total_placeholders": len(placeholders_data),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height
            }
        }

    def _analyze_slides(self, prs: Presentation) -> List[Dict[str, Any]]:
        """Analyze each slide in the presentation."""
        slides_data = []

        for idx, slide in enumerate(prs.slides):
            slide_info = {
                "slide_index": idx,
                "slide_id": slide.slide_id,
                "layout_name": (
                    slide.slide_layout.name
                    if hasattr(slide.slide_layout, 'name')
                    else "Unknown"
                ),
                "shape_count": len(slide.shapes),
                "has_title": any(
                    shape.has_text_frame and
                    hasattr(shape, 'name') and
                    'title' in shape.name.lower()
                    for shape in slide.shapes
                )
            }
            slides_data.append(slide_info)

        return slides_data

    def _extract_placeholders(self, prs: Presentation) -> List[Dict[str, Any]]:
        """Extract all placeholders from presentation."""
        placeholders = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape can hold text
                if shape.has_text_frame:
                    placeholder_info = {
                        "placeholder_id": f"slide_{slide_idx}_shape_{shape_idx}",
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_name": shape.name,
                        "shape_type": str(shape.shape_type),
                        "placeholder_type": "text",
                        "current_text": shape.text if shape.has_text_frame else "",
                        "max_chars": self._estimate_max_chars(shape),
                        "formatting_info": self._extract_formatting(shape)
                    }
                    placeholders.append(placeholder_info)

                # Check for tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    placeholder_info = {
                        "placeholder_id": f"slide_{slide_idx}_shape_{shape_idx}",
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_name": shape.name,
                        "placeholder_type": "table",
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns)
                    }
                    placeholders.append(placeholder_info)

        return placeholders

    def _estimate_max_chars(self, shape) -> int:
        """Estimate maximum characters based on shape dimensions."""
        if not shape.has_text_frame:
            return 0

        # Rough estimation based on dimensions
        width_chars = shape.width // 70000  # Approximate chars per width unit
        height_lines = shape.height // 200000  # Approximate lines per height unit

        return int(width_chars * height_lines * 80)  # ~80 chars per line

    def _extract_formatting(self, shape) -> Dict[str, Any]:
        """Extract formatting information from shape."""
        if not shape.has_text_frame:
            return {}

        formatting_info: Dict[str, Any] = {
            "font_name": None,
            "font_size": None,
            "bold": False,
            "italic": False,
            "alignment": None
        }

        try:
            if shape.text_frame.paragraphs:
                first_para = shape.text_frame.paragraphs[0]
                if first_para.runs:
                    first_run = first_para.runs[0]
                    formatting_info["font_name"] = first_run.font.name
                    formatting_info["font_size"] = first_run.font.size
                    formatting_info["bold"] = first_run.font.bold
                    formatting_info["italic"] = first_run.font.italic
                formatting_info["alignment"] = (
                    str(first_para.alignment) if first_para.alignment else None
                )
        except Exception:
            # Keep default values if extraction fails
            pass

        return formatting_info
